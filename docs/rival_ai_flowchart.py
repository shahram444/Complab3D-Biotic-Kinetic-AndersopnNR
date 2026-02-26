#!/usr/bin/env python3
"""Generate PowerPoint slides: Rival AI Flowchart + Code Reference.
   v2 — audience-friendly language, delta explained, clearer flow."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──────────────────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0E, 0x1A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x88, 0x88, 0x99)
LGRAY     = RGBColor(0xAA, 0xAA, 0xBB)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
TEAL      = RGBColor(0x2A, 0xCF, 0xAF)
RED       = RGBColor(0xE0, 0x60, 0x60)
RED_DARK  = RGBColor(0xB0, 0x30, 0x30)
BLUE      = RGBColor(0x4F, 0xA4, 0xFF)
GREEN     = RGBColor(0x4F, 0xDF, 0x6F)
ORANGE    = RGBColor(0xEF, 0x8F, 0x3F)
PURPLE    = RGBColor(0xCF, 0x6F, 0xFF)
DARK_BOX  = RGBColor(0x14, 0x18, 0x28)
BORDER    = RGBColor(0x2A, 0x2E, 0x42)
CODE_BG   = RGBColor(0x10, 0x14, 0x22)
CYAN      = RGBColor(0x50, 0xB0, 0xC0)
YELLOW    = RGBColor(0xFF, 0xFF, 0x5F)
COMMENT   = RGBColor(0x70, 0x88, 0x58)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, w, h, fill_color, border_color=None, border_w=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_diamond(slide, left, top, w, h, fill_color, border_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    shape.shadow.inherit = False
    return shape


def set_text(shape, text, size=12, color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Consolas"
    return tf


def add_text_multi(shape, lines, default_size=11, default_color=WHITE):
    """lines = [(text, size, color, bold, align), ...]"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        sz = line_data[1] if len(line_data) > 1 else default_size
        col = line_data[2] if len(line_data) > 2 else default_color
        bld = line_data[3] if len(line_data) > 3 else False
        al = line_data[4] if len(line_data) > 4 else PP_ALIGN.CENTER
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = al
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = "Consolas"


def add_line(slide, x1, y1, x2, y2, color=GRAY, width=Pt(2)):
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shape.line.color.rgb = color
    shape.line.width = width
    return shape


def add_free_text(slide, left, top, w, h, text, size=11, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Consolas"
    return txBox


def add_multiline_textbox(slide, left, top, w, h, lines):
    """lines = [(text, size, color, bold, align), ...]"""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        sz = line_data[1] if len(line_data) > 1 else 11
        col = line_data[2] if len(line_data) > 2 else WHITE
        bld = line_data[3] if len(line_data) > 3 else False
        al = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = al
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = "Consolas"
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: FLOWCHART  (audience-friendly)
# ══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide1, BG)

# Title
add_free_text(slide1, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
              "Rival Microbe AI — How It Decides What To Do", 28, GOLD, True, PP_ALIGN.CENTER)
add_free_text(slide1, Inches(0.5), Inches(0.62), Inches(12), Inches(0.35),
              "Every 0.4 seconds, the rival asks itself: \"What should I do next?\"",
              14, LGRAY, False, PP_ALIGN.CENTER)

# ── Flowchart layout ──────────────────────────────────────────────────────
cx = Inches(5.8)
dw = Inches(2.8)
dh = Inches(1.0)

y0 = Inches(1.15)
y1 = Inches(2.0)
y2 = Inches(3.05)
y3 = Inches(4.3)
y4 = Inches(5.5)
y5 = Inches(6.45)

ax_left = Inches(1.0)
ax_right = Inches(9.8)

# ── STEP 1: Timer fires ──
step_label = add_free_text(slide1, Inches(0.15), y0 + Inches(0.05), Inches(0.9), Inches(0.35),
                           "STEP 1", 9, GOLD, True, PP_ALIGN.CENTER)
start = add_box(slide1, cx - Inches(1.2), y0, Inches(2.4), Inches(0.6), RED_DARK, RED, Pt(2))
add_text_multi(start, [
    ("Timer fires", 13, WHITE, True),
    ("do_move() is called", 9, LGRAY),
])

add_line(slide1, cx, y0 + Inches(0.6), cx, y1, GRAY)

# ── STEP 2: Roll the dice ──
add_free_text(slide1, Inches(0.15), y1 + Inches(0.05), Inches(0.9), Inches(0.35),
              "STEP 2", 9, GOLD, True, PP_ALIGN.CENTER)
roll = add_box(slide1, cx - Inches(1.5), y1, Inches(3.0), Inches(0.65), DARK_BOX, BORDER)
add_text_multi(roll, [
    ("Roll a random number 0–100%", 12, WHITE, True),
    ("var roll = randf()  — like rolling a die", 9, LGRAY),
])

add_line(slide1, cx, y1 + Inches(0.65), cx, y2, GRAY)

# ── STEP 3 / DECISION 1: SEEK FOOD ──
add_free_text(slide1, Inches(0.15), y2 + Inches(0.15), Inches(0.9), Inches(0.35),
              "STEP 3", 9, GOLD, True, PP_ALIGN.CENTER)
d1 = add_diamond(slide1, cx - Inches(1.4), y2, dw, dh, RGBColor(0x0F, 0x15, 0x25), GREEN)
add_text_multi(d1, [
    ("Am I hungry?", 12, GREEN, True),
    ("roll < seek_chance", 9, LGRAY),
    ("(45% to 75%)", 9, LGRAY),
])

# YES → SEEK FOOD
add_line(slide1, cx - Inches(1.4), y2 + Inches(0.5), ax_left + Inches(2.8), y2 + Inches(0.5), GREEN)
seek_box = add_box(slide1, ax_left, y2 - Inches(0.15), Inches(2.8), Inches(1.3), RGBColor(0x0A, 0x20, 0x10), GREEN)
add_text_multi(seek_box, [
    ("SEEK FOOD", 15, GREEN, True),
    ("", 4, GREEN),
    ("Look around for the nearest", 10, WHITE),
    ("food (substrate) within", 10, WHITE),
    ("12 tiles — then move", 10, WHITE),
    ("one step toward it", 10, WHITE),
])
add_free_text(slide1, cx - Inches(2.3), y2 + Inches(0.15), Inches(0.7), Inches(0.3),
              "YES", 11, GREEN, True, PP_ALIGN.CENTER)

# NO → next
add_line(slide1, cx, y2 + dh, cx, y3, GRAY)
add_free_text(slide1, cx + Inches(0.1), y2 + dh - Inches(0.05), Inches(0.5), Inches(0.3),
              "NO", 11, RED, True)

# ── STEP 4 / DECISION 2: TRACK ARKE ──
add_free_text(slide1, Inches(0.15), y3 + Inches(0.15), Inches(0.9), Inches(0.35),
              "STEP 4", 9, GOLD, True, PP_ALIGN.CENTER)
d2 = add_diamond(slide1, cx - Inches(1.4), y3, dw, dh, RGBColor(0x0F, 0x15, 0x25), TEAL)
add_text_multi(d2, [
    ("Is the player nearby?", 12, TEAL, True),
    ("extra 15% chance &", 9, LGRAY),
    ("ARKE within 10 tiles", 9, LGRAY),
])

# YES → TRACK ARKE
add_line(slide1, cx + Inches(1.4), y3 + Inches(0.5), ax_right, y3 + Inches(0.5), TEAL)
track_box = add_box(slide1, ax_right, y3 - Inches(0.15), Inches(2.8), Inches(1.3), RGBColor(0x0A, 0x18, 0x18), TEAL)
add_text_multi(track_box, [
    ("TRACK ARKE", 15, TEAL, True),
    ("", 4, TEAL),
    ("The player (ARKE) is close!", 10, WHITE),
    ("Chase them to compete", 10, WHITE),
    ("for the same food zone.", 10, WHITE),
    ("Move one step toward ARKE.", 10, WHITE),
])
add_free_text(slide1, cx + Inches(1.5), y3 + Inches(0.15), Inches(0.7), Inches(0.3),
              "YES", 11, TEAL, True, PP_ALIGN.CENTER)

# NO → next
add_line(slide1, cx, y3 + dh, cx, y4, GRAY)
add_free_text(slide1, cx + Inches(0.1), y3 + dh - Inches(0.05), Inches(0.5), Inches(0.3),
              "NO", 11, RED, True)

# ── STEP 5 / DECISION 3: RIDE FLOW ──
add_free_text(slide1, Inches(0.15), y4 + Inches(0.15), Inches(0.9), Inches(0.35),
              "STEP 5", 9, GOLD, True, PP_ALIGN.CENTER)
d3 = add_diamond(slide1, cx - Inches(1.4), y4, dw, dh, RGBColor(0x0F, 0x15, 0x25), BLUE)
add_text_multi(d3, [
    ("Is there water flow?", 12, BLUE, True),
    ("flow speed > 0.1 &", 9, LGRAY),
    ("40% random chance", 9, LGRAY),
])

# YES → RIDE FLOW
add_line(slide1, cx - Inches(1.4), y4 + Inches(0.5), ax_left + Inches(2.8), y4 + Inches(0.5), BLUE)
flow_box = add_box(slide1, ax_left, y4 - Inches(0.15), Inches(2.8), Inches(1.3), RGBColor(0x0A, 0x12, 0x22), BLUE)
add_text_multi(flow_box, [
    ("RIDE THE FLOW", 15, BLUE, True),
    ("", 4, BLUE),
    ("There is a water current", 10, WHITE),
    ("at this tile. Drift along", 10, WHITE),
    ("with it (like a river).", 10, WHITE),
    ("Only if flow speed > 0.1", 10, WHITE),
])
add_free_text(slide1, cx - Inches(2.3), y4 + Inches(0.15), Inches(0.7), Inches(0.3),
              "YES", 11, BLUE, True, PP_ALIGN.CENTER)

# NO → WANDER
add_line(slide1, cx, y4 + dh, cx, y5, GRAY)
add_free_text(slide1, cx + Inches(0.1), y4 + dh - Inches(0.05), Inches(0.5), Inches(0.3),
              "NO", 11, RED, True)

# ── STEP 6: WANDER (fallback) ──
add_free_text(slide1, Inches(0.15), y5 + Inches(0.05), Inches(0.9), Inches(0.35),
              "STEP 6", 9, GOLD, True, PP_ALIGN.CENTER)
wander = add_box(slide1, cx - Inches(1.5), y5, Inches(3.0), Inches(0.7), RGBColor(0x20, 0x15, 0x0A), ORANGE)
add_text_multi(wander, [
    ("WANDER RANDOMLY", 14, ORANGE, True),
    ("Nothing else worked — pick a random direction", 9, LGRAY),
])

# ── RIGHT SIDE: Hunger mechanic note ──
hunger_box = add_box(slide1, Inches(10.0), y2 - Inches(0.6), Inches(3.0), Inches(2.0),
                      RGBColor(0x1A, 0x10, 0x10), RED)
add_text_multi(hunger_box, [
    ("HUNGER SYSTEM", 14, RED, True),
    ("", 4, GRAY),
    ("Every frame, hunger grows:", 10, WHITE),
    ("hunger += delta * 0.3", 10, CYAN),
    ("", 4, GRAY),
    ("The hungrier the rival,", 10, WHITE),
    ("the MORE LIKELY it will", 10, WHITE),
    ("choose \"Seek Food\".", 10, WHITE),
    ("", 4, GRAY),
    ("seek_chance:", 10, LGRAY),
    ("  low hunger  = 45%", 10, GREEN),
    ("  high hunger = 75%", 10, RED),
    ("", 4, GRAY),
    ("Eating resets hunger to 0", 9, GREEN),
])

# ── RIGHT SIDE: WHAT IS DELTA? ──
delta_box = add_box(slide1, Inches(10.0), y3 + Inches(0.1), Inches(3.0), Inches(2.1),
                     RGBColor(0x10, 0x10, 0x20), PURPLE)
add_text_multi(delta_box, [
    ("WHAT IS \"delta\"?", 14, PURPLE, True),
    ("", 4, GRAY),
    ("delta = the time (in seconds)", 10, WHITE),
    ("since the last frame was drawn.", 10, WHITE),
    ("", 4, GRAY),
    ("At 60 FPS:", 10, LGRAY),
    ("  delta = ~0.016 sec (1/60)", 10, CYAN),
    ("At 30 FPS:", 10, LGRAY),
    ("  delta = ~0.033 sec (1/30)", 10, CYAN),
    ("", 4, GRAY),
    ("Why use it? So the game runs", 10, WHITE),
    ("at the SAME SPEED on fast &", 10, WHITE),
    ("slow computers. Without delta,", 10, WHITE),
    ("a faster PC = faster hunger.", 10, WHITE),
    ("", 4, GRAY),
    ("hunger += delta * 0.3 means:", 9, GOLD),
    ("\"grow hunger by 0.3 per second,", 9, LGRAY),
    (" regardless of frame rate.\"", 9, LGRAY),
])

# ── Footer ──
add_free_text(slide1, Inches(0.3), Inches(7.15), Inches(12), Inches(0.3),
              "ARKE: Guardians of Earth  |  rival.gd  |  do_move() is the brain of every rival microbe",
              10, GRAY, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: CODE REFERENCE WITH INLINE COMMENTS (audience-friendly)
# ══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, BG)

# Title
add_free_text(slide2, Inches(0.5), Inches(0.1), Inches(12), Inches(0.5),
              "Rival AI — Code Walkthrough (Line by Line)", 26, GOLD, True, PP_ALIGN.CENTER)
add_free_text(slide2, Inches(0.5), Inches(0.55), Inches(12), Inches(0.3),
              "rival.gd  —  Each line is annotated so you can follow along",
              13, LGRAY, False, PP_ALIGN.CENTER)

# Code sections — rewritten with audience-friendly comments
code_sections = [
    {
        "title": "1. SETUP — Constants & Variables",
        "subtitle": "These values are set once when the rival spawns",
        "color": PURPLE,
        "code": [
            ("const SENSE_RANGE := 12",
             "# The rival can 'smell' food up to 12 tiles away"),
            ("const PLAYER_SENSE := 10",
             "# It can detect the player (ARKE) within 10 tiles"),
            ("var hunger: float = 0.0",
             "# Starts at 0; increases every frame over time"),
            ("var move_dir: int = 1",
             "# Its preferred wander direction (1=N, 2=E, 3=S, 4=W)"),
        ],
        "line_ref": "Lines 4-17"
    },
    {
        "title": "2. HUNGER GROWS — Called every frame",
        "subtitle": "delta makes this frame-rate independent (see Slide 1)",
        "color": RED,
        "code": [
            ("func _process(delta):",
             "# Godot calls this every single frame (~60x per sec)"),
            ("  hunger += delta * 0.3",
             "# delta = time since last frame (e.g. 0.016s at 60fps)"),
            ("  # So hunger grows by exactly 0.3 per second,",
             ""),
            ("  # no matter if the game runs at 30fps or 120fps.",
             ""),
        ],
        "line_ref": "Line 28-30"
    },
    {
        "title": "3. SEEK FOOD — Highest priority (45-75% chance)",
        "subtitle": "The hungrier the rival, the more likely it picks this",
        "color": GREEN,
        "code": [
            ("var seek_chance = clampf(",
             "# Calculate the chance to seek food:"),
            ("  0.45 + hunger * 0.05,",
             "#   Base 45% + 5% per hunger point"),
            ("  0.45, 0.75)",
             "#   But never below 45% or above 75%"),
            ("if roll < seek_chance:",
             "# Did we roll low enough? → Go find food!"),
            ("  var best = _seek_substrate(subs)",
             "# Scan all food sources on the map"),
            ("  for s in subs:",
             "# Loop through each food item one by one"),
            ("    if dist < SENSE_RANGE * T:",
             "# Is this food within our 12-tile smell range?"),
            ("      best_pos = s.position",
             "# Yes — remember it as the closest food"),
            ("  return _step_toward(best_pos)",
             "# Move one tile closer to that food"),
        ],
        "line_ref": "Lines 53-64, 117-132"
    },
    {
        "title": "4. TRACK ARKE — Chase the player (15% chance)",
        "subtitle": "Only happens if the player is close enough to notice",
        "color": TEAL,
        "code": [
            ("elif roll < seek_chance + 0.15:",
             "# Extra 15% chance to chase the player instead"),
            ("  var dx = player.tile_x - tile_x",
             "# How far away is ARKE horizontally?"),
            ("  var dy = player.tile_y - tile_y",
             "# How far away is ARKE vertically?"),
            ("  if absi(dx)+absi(dy) < PLAYER_SENSE:",
             "# Total distance < 10 tiles? Player is detected!"),
            ("    _step_toward(dx, dy)",
             "# Move one step toward ARKE's position"),
        ],
        "line_ref": "Lines 67-80"
    },
    {
        "title": "5. RIDE FLOW — Drift with the current (40% chance)",
        "subtitle": "Only available if the tile has a water current",
        "color": BLUE,
        "code": [
            ("var flow = world_ref.get_flow(",
             "# Ask the world: 'Is there a current at my tile?'"),
            ("  tile_x, tile_y)",
             "# Returns {dir: N/S/E/W, speed: 0.0-1.0}"),
            ("if flow['speed'] > 0.1",
             "# Is the current strong enough? (> 10% speed)"),
            ("   and randf() < 0.4:",
             "# AND a 40% random chance — not always drifting"),
            ("  var fd = flow_dirs[flow['dir']]",
             "# Convert 'North' → Vector2(0, -1) etc."),
            ("  _move_to(tile_x+fd.x, tile_y+fd.y)",
             "# Move one tile in the flow direction"),
        ],
        "line_ref": "Lines 86-96"
    },
    {
        "title": "6. WANDER — Last resort if nothing else triggers",
        "subtitle": "The rival just picks a direction and walks",
        "color": ORANGE,
        "code": [
            ("var dirs = [N, E, S, W]",
             "# All 4 possible directions to walk"),
            ("var preferred = dirs[move_dir-1]",
             "# Start with its favorite direction"),
            ("if is_walkable_tile(preferred):",
             "# Is there solid ground that way?"),
            ("  _move_to(preferred)",
             "# Yes — walk there"),
            ("  if randf() < 0.25:",
             "# 25% chance to randomly change"),
            ("    move_dir = randi_range(1,4)",
             "#   its favorite direction (keeps it varied)"),
            ("else: dirs.shuffle(); try each",
             "# Blocked? Try random directions instead"),
        ],
        "line_ref": "Lines 98-115"
    },
    {
        "title": "7. EATING — What happens when the rival finds food",
        "subtitle": "Resets hunger so the rival calms down",
        "color": YELLOW,
        "code": [
            ("func on_eat() -> void:",
             "# This is called when the rival reaches food"),
            ("  eat_flash = 0.3",
             "# Show a red flash for 0.3 sec (visual feedback)"),
            ("  hunger = 0.0",
             "# Reset hunger to 0 → seek_chance drops to 45%"),
            ("  # The rival is now 'full' and more likely",
             ""),
            ("  # to wander or ride flow instead of seeking.",
             ""),
        ],
        "line_ref": "Lines 167-171"
    },
]

# Layout: 2 columns
col_w = Inches(6.25)
col_gap = Inches(0.2)
left_x = Inches(0.4)
right_x = left_x + col_w + col_gap
start_y = Inches(0.95)
section_gap = Inches(0.06)

cur_y_left = start_y
cur_y_right = start_y

for idx, section in enumerate(code_sections):
    # Left column: sections 0-3, Right column: sections 4-6
    if idx < 4:
        cur_x = left_x
        cur_y = cur_y_left
    else:
        cur_x = right_x
        cur_y = cur_y_right

    # Section title bar
    title_h = Inches(0.42)
    title_box = add_box(slide2, cur_x, cur_y, col_w, title_h,
                        RGBColor(0x0E, 0x12, 0x20), section["color"], Pt(1.5))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(8)

    # Title line
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {section['title']}"
    run.font.size = Pt(11)
    run.font.color.rgb = section["color"]
    run.font.bold = True
    run.font.name = "Consolas"
    run2 = p.add_run()
    run2.text = f"    ({section['line_ref']})"
    run2.font.size = Pt(8)
    run2.font.color.rgb = GRAY
    run2.font.name = "Consolas"

    # Subtitle line
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(0)
    run3 = p2.add_run()
    run3.text = f"    {section['subtitle']}"
    run3.font.size = Pt(9)
    run3.font.color.rgb = LGRAY
    run3.font.name = "Consolas"
    run3.font.italic = True

    cur_y += title_h

    # Code lines
    code_h_per_line = Inches(0.20)
    code_h = code_h_per_line * len(section["code"]) + Inches(0.1)
    code_box = add_box(slide2, cur_x, cur_y, col_w, code_h,
                       CODE_BG, RGBColor(0x1A, 0x1E, 0x30), Pt(1))

    tf = code_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(4)

    for ci, (code, comment) in enumerate(section["code"]):
        if ci == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        p.space_before = Pt(1)

        # Code part
        run = p.add_run()
        run.text = code
        run.font.size = Pt(9)
        run.font.color.rgb = CYAN
        run.font.name = "Consolas"

        # Comment part
        if comment:
            run2 = p.add_run()
            run2.text = f"  {comment}"
            run2.font.size = Pt(8.5)
            run2.font.color.rgb = COMMENT
            run2.font.name = "Consolas"

    cur_y += code_h + section_gap

    if idx < 4:
        cur_y_left = cur_y
    else:
        cur_y_right = cur_y

# Footer
add_free_text(slide2, Inches(0.3), Inches(7.1), Inches(12), Inches(0.3),
              "ARKE: Guardians of Earth  |  rival.gd  |  GDScript (Godot 4.x)  |  Read each green comment left-to-right",
              10, GRAY, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: PRIORITY SUMMARY (clean audience-friendly overview)
# ══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3, BG)

add_free_text(slide3, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
              "Rival AI — The 4 Behaviors at a Glance", 28, GOLD, True, PP_ALIGN.CENTER)
add_free_text(slide3, Inches(0.5), Inches(0.6), Inches(12), Inches(0.3),
              "The rival checks these in order. The first one that passes wins.",
              14, LGRAY, False, PP_ALIGN.CENTER)

# Priority cascade boxes
priorities = [
    ("1", "SEEK\nFOOD", "45-75%",
     "The rival looks for the\nclosest food within 12\ntiles and walks toward it.\n\nHigher hunger = higher\nchance (45% up to 75%).",
     GREEN),
    ("2", "TRACK\nARKE", "+15%",
     "If the player is within\n10 tiles, the rival may\nchase them to compete\nfor the same food area.\n\nOnly a 15% extra chance.",
     TEAL),
    ("3", "RIDE\nFLOW", "40%*",
     "If there is a water\ncurrent at this tile\n(speed > 0.1), the rival\ndrifts with it.\n\n40% chance per decision.",
     BLUE),
    ("4", "WANDER", "Always",
     "Nothing else triggered.\nPick a direction and walk.\nTry preferred dir first;\nif blocked, shuffle and\ntry any open direction.\n25% chance to change dir.",
     ORANGE),
]

box_w = Inches(2.7)
box_h = Inches(3.6)
total_w = box_w * 4 + Inches(0.4) * 3
start_x = (Inches(13.333) - total_w) / 2
box_y = Inches(1.3)

for i, (num, title, chance, desc, color) in enumerate(priorities):
    bx = int(start_x) + i * int(box_w + Inches(0.4))

    # Main box
    box = add_box(slide3, bx, box_y, box_w, box_h, DARK_BOX, color, Pt(2))

    # Priority number circle
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL,
        bx + int(box_w/2) - Inches(0.3), box_y - Inches(0.25),
        Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    set_text(circle, num, 20, BG, True)

    # Content
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_top = Pt(35)
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)

    # Title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.color.rgb = color
    run.font.bold = True
    run.font.name = "Consolas"

    # Chance label
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    run2 = p2.add_run()
    run2.text = f"Chance: {chance}"
    run2.font.size = Pt(12)
    run2.font.color.rgb = GOLD
    run2.font.bold = True
    run2.font.name = "Consolas"

    # Description (plain English)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(14)
    run3 = p3.add_run()
    run3.text = desc
    run3.font.size = Pt(10)
    run3.font.color.rgb = LGRAY
    run3.font.name = "Consolas"

    # Arrow between boxes (except last)
    if i < 3:
        arrow_x = bx + box_w + Inches(0.02)
        arrow_y = box_y + int(box_h / 2)
        add_free_text(slide3, arrow_x, arrow_y - Inches(0.15), Inches(0.35), Inches(0.35),
                      ">>", 16, GRAY, True, PP_ALIGN.CENTER)

# Footnote
add_free_text(slide3, Inches(0.5), Inches(5.3), Inches(12), Inches(0.4),
              "* Ride Flow only happens when the tile actually has a water current (speed > 0.1)",
              11, GRAY, False, PP_ALIGN.CENTER)

# Hunger + Delta explanation at bottom
bottom_box = add_box(slide3, Inches(1.5), Inches(5.7), Inches(10.3), Inches(1.4),
                      RGBColor(0x12, 0x0E, 0x18), PURPLE, Pt(1.5))
add_text_multi(bottom_box, [
    ("KEY CONCEPTS", 14, PURPLE, True),
    ("", 3, GRAY),
    ("HUNGER:  hunger += delta * 0.3  means hunger grows by 0.3 every second.", 11, WHITE),
    ("         The hungrier the rival, the more it prioritizes finding food (45% -> 75%).", 11, LGRAY),
    ("         Eating food resets hunger back to 0 — the rival \"calms down.\"", 11, GREEN),
    ("", 3, GRAY),
    ("DELTA:   delta = seconds since the last frame.  At 60 FPS that is ~0.016 seconds.", 11, WHITE),
    ("         Multiplying by delta ensures the game behaves identically on all hardware.", 11, LGRAY),
    ("         Without it, a faster computer would make the rival get hungry faster!", 11, RED),
])

# Footer
add_free_text(slide3, Inches(0.3), Inches(7.15), Inches(12), Inches(0.3),
              "ARKE: Guardians of Earth  |  rival.gd  |  GDScript (Godot 4.x)  |  Based on CompLaB3D Research",
              10, GRAY, False, PP_ALIGN.CENTER)


# ── SAVE (overwrite existing) ─────────────────────────────────────────────────
output_path = "/home/user/Complab3D-Biotic-Kinetic-AndersopnNR/docs/Rival_AI_Flowchart.pptx"
prs.save(output_path)
print(f"Saved (overwritten): {output_path}")
